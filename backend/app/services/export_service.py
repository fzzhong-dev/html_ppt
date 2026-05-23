import json
import logging
import re
import subprocess
import tempfile
from pathlib import Path

from bs4 import BeautifulSoup
from pptx import Presentation as PptxPresentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from app.config import settings

logger = logging.getLogger(__name__)


class ExportService:
    def __init__(self):
        self.screenshot_script = Path(settings.screenshot_script)
        self.screenshots_dir = Path(settings.screenshots_dir)

    async def export_to_pptx(self, slides: list, title: str = "Presentation") -> str:
        """Export slides to an editable PPTX with native text boxes."""
        if not slides:
            raise ValueError("No slides to export")

        output_dir = self.screenshots_dir
        output_dir.mkdir(parents=True, exist_ok=True)
        pptx_path = str(output_dir / f"{title}.pptx")
        self.create_editable_pptx(slides, title, pptx_path)
        return pptx_path

    def create_editable_pptx(self, slides: list, title: str, output_path: str) -> str:
        prs = PptxPresentation()
        # 16:9 aspect ratio
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        blank_layout = prs.slide_layouts[6]  # blank layout

        for slide_data in slides:
            html = slide_data.html_content if hasattr(slide_data, 'html_content') else slide_data.get('html_content', '')
            slide = prs.slides.add_slide(blank_layout)
            self._populate_slide_from_html(slide, html, prs)

        prs.save(output_path)
        return output_path

    def _populate_slide_from_html(self, pptx_slide, html: str, prs) -> None:
        """Parse HTML and populate a PPTX slide with native text boxes."""
        try:
            soup = BeautifulSoup(html, 'html.parser')
        except Exception:
            logger.warning("Failed to parse HTML for export", exc_info=True)
            return

        body = soup.find('body')
        if not body:
            return

        # Set slide background from body style
        bg_color = self._extract_bg_color(body)
        if bg_color:
            self._set_slide_bg(pptx_slide, bg_color)

        # Extract text blocks from the HTML
        text_blocks = self._extract_text_blocks(body)

        slide_w = 13.333  # inches
        slide_h = 7.5

        for block in text_blocks:
            text = block['text'].strip()
            if not text:
                continue

            # Position: estimate from CSS or use auto-layout
            left = block.get('left', 0.5)
            top = block.get('top', 0.5)
            width = block.get('width', slide_w - 1.0)
            height = block.get('height', 0)

            # Font size estimation
            font_size = block.get('font_size', 18)
            bold = block.get('bold', False)
            color = block.get('color', None)
            alignment = block.get('alignment', PP_ALIGN.LEFT)

            txBox = pptx_slide.shapes.add_textbox(
                Inches(left), Inches(top),
                Inches(width), Inches(max(height, 0.3)),
            )
            tf = txBox.text_frame
            tf.word_wrap = True

            # Split text into paragraphs
            paragraphs = text.split('\n')
            for pi, para_text in enumerate(paragraphs):
                para_text = para_text.strip()
                if not para_text:
                    continue
                if pi == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()

                run = p.add_run()
                run.text = para_text
                run.font.size = Pt(font_size)
                if bold:
                    run.font.bold = True
                if color:
                    rgb = self._hex_to_rgb(color)
                    if rgb:
                        run.font.color.rgb = rgb
                p.alignment = alignment
                p.space_after = Pt(4)

    def _extract_text_blocks(self, body) -> list[dict]:
        """Extract structured text blocks from HTML body."""
        blocks = []

        # Find elements with visible text
        for el in body.find_all(['h1', 'h2', 'h3', 'h4', 'h5', 'h6', 'p', 'li', 'span', 'div', 'td']):
            # Skip elements that contain other block elements (only get leaf-like text nodes)
            text = self._get_direct_text(el)
            if not text or len(text.strip()) < 2:
                continue

            style_str = self._get_element_style(el)
            block = {
                'text': text.strip(),
                'font_size': self._parse_font_size(style_str, 18),
                'bold': self._is_bold(el, style_str),
                'color': self._parse_color(style_str),
                'alignment': self._parse_alignment(style_str),
                'left': 0.5,
                'top': 0,
                'width': 12.3,
                'height': 0,
            }

            # Scale font size for PPTX (HTML 1920px → PPTX 13.33in)
            tag = el.name
            if tag in ('h1',):
                block['font_size'] = max(block['font_size'], 36)
                block['bold'] = True
            elif tag in ('h2',):
                block['font_size'] = max(block['font_size'], 28)
                block['bold'] = True
            elif tag in ('h3', 'h4'):
                block['font_size'] = max(block['font_size'], 22)
                block['bold'] = True

            # Skip duplicate text (nested elements)
            text_clean = text.strip()
            if any(text_clean in b['text'] and b['text'] != text_clean for b in blocks):
                continue

            blocks.append(block)

        # Auto-layout: distribute blocks vertically if no position info
        if blocks:
            total_estimated = sum(max(0.3, len(b['text']) / 80 * 0.4) for b in blocks)
            available_h = 7.0  # inches, minus margins
            if total_estimated > 0:
                y = 0.4
                for b in blocks:
                    h = max(0.3, len(b['text']) / 80 * 0.4)
                    b['top'] = y
                    b['height'] = h
                    y += h + 0.1
                    if y > available_h:
                        break

        return blocks

    def _get_direct_text(self, el) -> str:
        """Get text from element, preferring direct text over recursive."""
        # Get direct text nodes
        parts = []
        for child in el.children:
            if isinstance(child, str):
                parts.append(child.strip())
            elif child.name in ('span', 'strong', 'em', 'b', 'i', 'a', 'mark', 'small'):
                parts.append(child.get_text(strip=True))
        direct = ' '.join(p for p in parts if p)
        if direct:
            return direct
        return el.get_text(strip=True)

    def _get_element_style(self, el) -> str:
        style = el.get('style', '')
        # Also check inline <style> blocks and class-based styles
        return style

    def _parse_font_size(self, style_str: str, default: int = 18) -> int:
        m = re.search(r'font-size\s*:\s*(\d+)(?:px)?', style_str)
        if m:
            px = int(m.group(1))
            # Scale from 1920px canvas to ~13.33in PPTX
            return max(10, int(px * 0.75))
        return default

    def _is_bold(self, el, style_str: str) -> bool:
        if el.name in ('h1', 'h2', 'h3', 'h4', 'h5', 'h6', 'b', 'strong', 'th'):
            return True
        if 'font-weight' in style_str:
            m = re.search(r'font-weight\s*:\s*(\d+|bold)', style_str)
            if m:
                val = m.group(1)
                return val == 'bold' or (val.isdigit() and int(val) >= 600)
        return False

    def _parse_color(self, style_str: str) -> str | None:
        m = re.search(r'color\s*:\s*(#[0-9a-fA-F]{3,8})', style_str)
        if m:
            return m.group(1)
        m = re.search(r'color\s*:\s*rgb\((\d+)\s*,\s*(\d+)\s*,\s*(\d+)\)', style_str)
        if m:
            return f'#{int(m.group(1)):02x}{int(m.group(2)):02x}{int(m.group(3)):02x}'
        return None

    def _parse_alignment(self, style_str: str) -> int:
        if 'text-align' in style_str:
            if 'center' in style_str:
                return PP_ALIGN.CENTER
            if 'right' in style_str:
                return PP_ALIGN.RIGHT
        return PP_ALIGN.LEFT

    def _extract_bg_color(self, body) -> str | None:
        style = body.get('style', '')
        m = re.search(r'background(?:-color)?\s*:\s*(#[0-9a-fA-F]{3,8})', style)
        if m:
            return m.group(1)
        # Check for linear-gradient
        m = re.search(r'background\s*:\s*linear-gradient\([^,]+,\s*(#[0-9a-fA-F]{6})', style)
        if m:
            return m.group(1)
        return None

    def _set_slide_bg(self, pptx_slide, hex_color: str) -> None:
        rgb = self._hex_to_rgb(hex_color)
        if not rgb:
            return
        try:
            background = pptx_slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = rgb
        except Exception:
            logger.warning("Failed to set slide background", exc_info=True)

    def _hex_to_rgb(self, hex_color: str) -> RGBColor | None:
        hex_color = hex_color.lstrip('#')
        if len(hex_color) == 3:
            hex_color = ''.join(c * 2 for c in hex_color)
        if len(hex_color) < 6:
            return None
        try:
            return RGBColor(
                int(hex_color[0:2], 16),
                int(hex_color[2:4], 16),
                int(hex_color[4:6], 16),
            )
        except (ValueError, IndexError):
            return None

    # --- Legacy screenshot-based export (fallback) ---

    async def export_to_pptx_screenshots(self, slides: list, title: str = "Presentation") -> str:
        if not slides:
            raise ValueError("No slides to export")

        slides_data = {
            "slides": [
                {"page_number": s.page_number, "html_content": s.html_content}
                for s in slides
            ]
        }

        with tempfile.NamedTemporaryFile(mode='w', suffix='.json', delete=False, encoding='utf-8') as f:
            json.dump(slides_data, f, ensure_ascii=False)
            input_path = f.name

        output_dir = self.screenshots_dir / f"export_{id(slides)}"
        output_dir.mkdir(parents=True, exist_ok=True)

        result = subprocess.run(
            ["node", str(self.screenshot_script), input_path, str(output_dir)],
            capture_output=True, text=True, timeout=120,
        )

        if result.returncode != 0:
            raise RuntimeError(f"Screenshot failed: {result.stderr}")

        screenshots = json.loads(result.stdout)

        pptx_path = str(self.screenshots_dir / f"{title}.pptx")
        self.create_pptx_from_slides(screenshots, title, pptx_path)

        Path(input_path).unlink(missing_ok=True)

        return pptx_path

    def create_pptx_from_slides(self, screenshots: list[dict], title: str, output_path: str = None) -> str:
        if not screenshots:
            raise ValueError("No screenshots to export")

        prs = PptxPresentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(5.63)

        blank_layout = prs.slide_layouts[6]

        for shot in screenshots:
            slide = prs.slides.add_slide(blank_layout)
            slide.shapes.add_picture(
                shot["path"],
                Inches(0), Inches(0),
                Inches(10), Inches(5.63),
            )

        if output_path is None:
            output_path = str(self.screenshots_dir / f"{title}.pptx")
        prs.save(output_path)
        return output_path
